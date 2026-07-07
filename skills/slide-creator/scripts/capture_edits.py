# /// script
# requires-python = ">=3.9"
# dependencies = ["python-pptx"]
# ///
"""Capture manual PowerPoint edits by diffing an EDITED deck against a REFERENCE
deck regenerated from the current source.

This is the *untagged* fallback for round-tripping. When a build script was set
up round-trip-ready (scid object names + state.json), prefer sync_from_pptx.py
(see references/roundtrip.md) -- it maps each edited shape back to its exact
source object. Use this tool when the script is plain pptxgenjs with coordinate
constants and no scid names: it recovers "what did the human move/resize/retype
in PowerPoint" well enough to fold back into the coordinate constants by hand.

It never writes to any deck -- it only reads both files and prints a report.

Workflow
--------
1. Back up the edited file; never let a build overwrite it.
2. Confirm the edited file was actually saved (mtime newer than last build).
3. Regenerate a reference from the CURRENT source to a TEMP path, replaying the
   FULL post-build chain (add_equation.py / add_video.py with identical coords).
4. uv run scripts/capture_edits.py --reference REF.pptx --edited EDITED.pptx
5. Fold reported deltas into the build script; rebuild to temp; re-run until the
   residual is ~0.

Matching strategy (in order): exact text -> content signature (image sha1,
media poster sha1 + relationship targets, chart part, table text) -> positional
proximity. Content shapes are matched by *signature*, not by type name, so two
same-size images are told apart; duplicates that share a signature are paired by
position and reported as medium/low confidence. PowerPoint title-autofit
(width-only change on a title-like box) is flagged as suspected noise.
"""
import argparse
import hashlib
import json
import os
import sys
from collections import defaultdict, deque

try:
    from pptx import Presentation
    try:
        from pptx.exc import PackageNotFoundError
    except Exception:  # very old python-pptx
        class PackageNotFoundError(Exception):
            pass
except Exception as exc:  # pragma: no cover - import environment problem
    print(f"error: python-pptx import failed: {exc}", file=sys.stderr)
    sys.exit(2)

EMU = 914400  # EMU per inch

# Shapes that carry real content even without extractable text. A moved / added
# / removed one of these is a genuine edit -- but only if we can pin it to a
# content signature. Type name alone is NOT enough (two images are both PICTURE).
CONTENT_TYPES = {"PICTURE", "LINKED_PICTURE", "MEDIA", "CHART", "TABLE"}

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
}
R_EMBED = "{%s}embed" % NS["r"]
R_LINK = "{%s}link" % NS["r"]

TITLE_BAND_IN = 1.35   # a text box whose top is above this is "title-like"
TITLE_MIN_W_IN = 6.0   # ...and at least this wide


# --------------------------------------------------------------------------- #
# shape metadata helpers
# --------------------------------------------------------------------------- #
def _iter(shapes):
    """Flatten group shapes one level so children are comparable."""
    for sh in shapes:
        st = sh.shape_type
        if st is not None and str(st).startswith("GROUP"):
            try:
                yield from _iter(sh.shapes)
                continue
            except Exception:
                pass
        yield sh


def _text(sh):
    try:
        # collapse all whitespace so PowerPoint run-merging (extra/removed
        # spaces on save) does not read as a text edit.
        return " ".join((sh.text or "").split())
    except Exception:
        return ""


def _linelike(sh):
    w, h = sh.width, sh.height
    thin = 0.05 * EMU
    return (w is not None and w < thin) or (h is not None and h < thin)


def _type(sh):
    # has_table / has_chart are reliable; shape_type is None for tables.
    try:
        if getattr(sh, "has_table", False) and sh.has_table:
            return "TABLE"
    except Exception:
        pass
    try:
        if getattr(sh, "has_chart", False) and sh.has_chart:
            return "CHART"
    except Exception:
        pass
    st = sh.shape_type
    return str(st).split()[0] if st is not None else "UNKNOWN"


def _geo(sh):
    return (sh.left, sh.top, sh.width, sh.height, int(getattr(sh, "rotation", 0) or 0))


def _center(sh):
    if sh.left is None or sh.top is None:
        return None
    return (sh.left + (sh.width or 0) / 2, sh.top + (sh.height or 0) / 2)


def _in(v):
    return None if v is None else round(v / EMU, 3)


def _bbox(sh):
    l, t, w, h, r = _geo(sh)
    return {"x": _in(l), "y": _in(t), "w": _in(w), "h": _in(h), "rot": r}


def _title_like(sh):
    if not _text(sh):
        return False
    l, t, w, h, _ = _geo(sh)
    if t is None or w is None:
        return False
    return t <= TITLE_BAND_IN * EMU and w >= TITLE_MIN_W_IN * EMU


def _rid_blob_sha1(part, rid):
    try:
        rel = part.rels[rid]
        if rel.is_external:
            return None
        return hashlib.sha1(rel.target_part.blob).hexdigest()
    except Exception:
        return None


def _blip_sha1(sh):
    """sha1 of the first embedded raster blip (picture pixels or media poster)."""
    try:
        el = sh._element
        part = sh.part
    except Exception:
        return None
    for blip in el.findall(".//a:blip", NS):
        rid = blip.get(R_EMBED)
        if rid:
            s = _rid_blob_sha1(part, rid)
            if s:
                return s
    return None


def _media_targets(sh):
    """Relationship targets of the video/audio/media parts (filenames/URLs)."""
    try:
        el = sh._element
        part = sh.part
    except Exception:
        return []
    out = []
    for tag in ("a:videoFile", "a:audioFile", "p14:media"):
        for m in el.findall(".//" + tag, NS):
            rid = m.get(R_LINK) or m.get(R_EMBED)
            if not rid:
                continue
            try:
                rel = part.rels[rid]
                out.append(rel.target_ref if rel.is_external else str(rel.target_part.partname))
            except Exception:
                pass
    return sorted(set(out))


def _content_sig(sh, typ):
    """Return (signature, missing_reason). signature=None means no content
    signature could be extracted (caller downgrades confidence)."""
    if typ == "TABLE":
        try:
            txt = "|".join(c.text for row in sh.table.rows for c in row.cells)
            return "table:" + hashlib.sha1(txt.encode("utf-8")).hexdigest()[:16], None
        except Exception:
            return None, "table-unreadable"
    if typ == "CHART":
        try:
            return "chart:" + str(sh.chart.part.partname), None
        except Exception:
            return None, "chart-unreadable"
    if typ in ("PICTURE", "LINKED_PICTURE"):
        try:
            return "pic:" + sh.image.sha1, None
        except Exception:
            s = _blip_sha1(sh)
            return ("pic:" + s, None) if s else (None, "picture-no-embedded-blob")
    if typ == "MEDIA":
        poster = _blip_sha1(sh)
        tgts = _media_targets(sh)
        if poster or tgts:
            return "media:" + (poster or "-") + "|" + ";".join(tgts), None
        return None, "media-no-signature"
    return None, None  # non-content shape: signature not applicable


def _meta(sh):
    typ = _type(sh)
    sig, sig_reason = _content_sig(sh, typ)
    return {
        "sh": sh, "type": typ, "text": _text(sh), "sig": sig, "sig_reason": sig_reason,
        "center": _center(sh), "title_like": _title_like(sh),
    }


# --------------------------------------------------------------------------- #
# diffing
# --------------------------------------------------------------------------- #
def _deltas(a, b, th):
    out = []
    for name, x, y in zip("xywh", _geo(a)[:4], _geo(b)[:4]):
        if x is not None and y is not None and abs(x - y) > th:
            out.append({"prop": name, "from": _in(x), "to": _in(y)})
    if _geo(a)[4] != _geo(b)[4]:
        out.append({"prop": "rot", "from": _geo(a)[4], "to": _geo(b)[4]})
    return out


def _greedy_pairs(refs, eds, tol):
    """Pair ref->edited by nearest center within tol (same type). Returns
    (pairs, leftover_refs, leftover_eds). Each pair carries an 'ambiguous' flag
    if more than one edited candidate fell within tol."""
    used_e = set()
    matched_r = set()
    pairs = []
    for i, r in enumerate(refs):
        cr = r["center"]
        best, bd, cand = None, None, 0
        for j, e in enumerate(eds):
            if j in used_e or e["type"] != r["type"]:
                continue
            if _linelike(r["sh"]) != _linelike(e["sh"]):
                continue  # never pair a zero-width line to a solid box
            ce = e["center"]
            if cr is None or ce is None:
                continue
            d = abs(cr[0] - ce[0]) + abs(cr[1] - ce[1])
            if d <= tol:
                cand += 1
                if best is None or d < bd:
                    best, bd = j, d
        if best is not None:
            used_e.add(best)
            matched_r.add(i)
            pairs.append((r, eds[best], cand > 1))
    leftover_r = [r for i, r in enumerate(refs) if i not in matched_r]
    leftover_e = [e for j, e in enumerate(eds) if j not in used_e]
    return pairs, leftover_r, leftover_e


SUGGEST = {
    ("moved", "high"): "Fold the bbox change into the build-script coordinate constants.",
    ("moved", "medium"): "Fold the bbox change into the build script; verify the match first.",
    ("moved", "low"): "Verify manually before editing source (ambiguous / possible noise).",
    ("text", "high"): "Update the text (and any bbox change) in the build script.",
    ("text", "medium"): "Update the text in the build script; verify the match first.",
    ("added", "*"): "New shape added by hand -- recreate it in source (or import via round-trip).",
    ("removed", "medium"): "Shape removed in PowerPoint -- remove from source or confirm.",
    ("removed", "low"): "Likely equation/decorative re-serialization -- verify, probably ignore.",
    ("noise", "*"): "PowerPoint autofit (width-only on a title) -- safe to ignore.",
}


def _record(slide, kind, matched_by, confidence, reason, **kw):
    rec = {"slide": slide, "kind": kind, "matched_by": matched_by,
           "confidence": confidence, "reason": reason,
           "suspected_noise": kw.get("suspected_noise", False)}
    for k in ("before_bbox", "after_bbox", "text_from", "text_to", "deltas"):
        if k in kw:
            rec[k] = kw[k]
    grp = "moved"
    if kind in ("added_in_edited",):
        grp = "added"
    elif kind in ("removed_in_edited",):
        grp = "removed"
    elif kind == "text_changed":
        grp = "text"
    if rec["suspected_noise"]:
        rec["suggested_action"] = SUGGEST[("noise", "*")]
    else:
        rec["suggested_action"] = (SUGGEST.get((grp, confidence))
                                   or SUGGEST.get((grp, "*"))
                                   or "Review.")
    return rec


def _emit_pair(slide, r, e, ambiguous, matched_by, th):
    d = _deltas(r["sh"], e["sh"], th)
    text_changed = r["text"] != e["text"] and (r["text"] or e["text"])
    if not d and not text_changed:
        return None
    # confidence from match quality
    if matched_by == "text":
        conf = "high"
    elif matched_by == "content":
        conf = "medium" if ambiguous else "high"
    else:  # proximity
        conf = "low" if ambiguous else "medium"
    reason = {
        "text": "matched by exact text",
        "content": "matched by content signature" + (" (duplicate sig, paired by position)" if ambiguous else ""),
        "proximity": "matched by position" + (" (ambiguous: multiple candidates)" if ambiguous else "") + (" (text changed)" if text_changed else ""),
    }[matched_by]
    mb = "ambiguous" if (ambiguous and matched_by == "proximity") else matched_by
    # title autofit suppression: width-only change on a title-like box
    suspected = False
    if d and all(x["prop"] == "w" for x in d) and (r["title_like"] or e["title_like"]):
        suspected = True
        conf = "low"
        reason = "PowerPoint title autofit (width-only)"
    kind = "text_changed" if text_changed and not d else "moved"
    kw = {"before_bbox": _bbox(r["sh"]), "after_bbox": _bbox(e["sh"]),
          "deltas": d, "suspected_noise": suspected}
    if text_changed:
        kw["text_from"] = r["text"][:120]
        kw["text_to"] = e["text"][:120]
    return _record(slide, kind, mb, conf, reason, **kw)


def diff_slide(slide_no, ref_shapes, ed_shapes, th, tol):
    refs = [_meta(s) for s in _iter(ref_shapes)]
    eds = [_meta(s) for s in _iter(ed_shapes)]
    changes = []

    # Pass A: exact text
    ebt = defaultdict(deque)
    for e in eds:
        if e["text"]:
            ebt[e["text"]].append(e)
    rem_r, matched_e = [], set()
    for r in refs:
        if r["text"] and ebt[r["text"]]:
            e = ebt[r["text"]].popleft()
            matched_e.add(id(e))
            rec = _emit_pair(slide_no, r, e, False, "text", th)
            if rec:
                changes.append(rec)
        else:
            rem_r.append(r)
    rem_e = [e for e in eds if id(e) not in matched_e]

    # Pass B: content signature (distinguishes same-type shapes by pixels/media)
    ebs = defaultdict(list)
    for e in rem_e:
        if e["sig"]:
            ebs[e["sig"]].append(e)
    still_r = []
    used_e = set()
    for r in rem_r:
        if r["sig"] and ebs.get(r["sig"]):
            bucket = [e for e in ebs[r["sig"]] if id(e) not in used_e]
            if not bucket:
                still_r.append(r)
                continue
            dup = len(ebs[r["sig"]]) > 1
            if len(bucket) == 1:
                e = bucket[0]
            else:  # duplicate signatures -> pair by nearest position
                cr = r["center"]
                e = min(bucket, key=lambda x: (abs((x["center"] or (0, 0))[0] - (cr or (0, 0))[0])
                                               + abs((x["center"] or (0, 0))[1] - (cr or (0, 0))[1])))
            used_e.add(id(e))
            rec = _emit_pair(slide_no, r, e, dup, "content", th)
            if rec:
                changes.append(rec)
        else:
            still_r.append(r)
    rem_e2 = [e for e in rem_e if id(e) not in used_e]

    # Pass C: positional proximity (text edits + sig-less content moves)
    pairs, leftover_r, leftover_e = _greedy_pairs(still_r, rem_e2, tol)
    for r, e, amb in pairs:
        rec = _emit_pair(slide_no, r, e, amb, "proximity", th)
        if rec:
            changes.append(rec)

    # Unmatched
    for r in leftover_r:
        content = r["type"] in CONTENT_TYPES
        conf = "medium" if (r["text"] or (content and r["sig"])) else "low"
        reason = ("content shape (%s) gone; signature=%s" % (r["type"], "yes" if r["sig"] else "no")
                  if content else "text/shape not found in edited (may be re-serialization)")
        changes.append(_record(slide_no, "removed_in_edited", "unmatched", conf, reason,
                               before_bbox=_bbox(r["sh"]),
                               **({"text_from": r["text"][:120]} if r["text"] else {})))
    for e in leftover_e:
        content = e["type"] in CONTENT_TYPES
        conf = "medium" if (e["text"] or (content and e["sig"])) else "low"
        reason = ("new content shape (%s)" % e["type"] if content
                  else "new text/shape (may be re-serialization)")
        changes.append(_record(slide_no, "added_in_edited", "unmatched", conf, reason,
                               after_bbox=_bbox(e["sh"]),
                               **({"text_to": e["text"][:120]} if e["text"] else {})))
    return changes


# --------------------------------------------------------------------------- #
# CLI
# --------------------------------------------------------------------------- #
def _load(path, label):
    if not os.path.isfile(path):
        print(f"error: --{label} file not found: {path}", file=sys.stderr)
        sys.exit(2)
    try:
        return Presentation(path)
    except PackageNotFoundError:
        print(f"error: --{label} is not a valid .pptx package: {path}", file=sys.stderr)
        sys.exit(2)
    except Exception as exc:
        print(f"error: failed to open --{label} ({path}): {exc}", file=sys.stderr)
        sys.exit(2)


def _parse_slides(spec):
    if not spec:
        return None
    out = set()
    for tok in spec.replace(" ", "").split(","):
        if not tok:
            continue
        if not tok.isdigit():
            print(f"error: --slide takes comma-separated slide numbers, got '{tok}'", file=sys.stderr)
            sys.exit(2)
        out.add(int(tok))
    return out


def main():
    ap = argparse.ArgumentParser(
        description="Diff a hand-edited PPTX against a reference regenerated from the current source. Read-only; never writes a deck.")
    ap.add_argument("--reference", required=True, help="deck rebuilt from current source (temp path)")
    ap.add_argument("--edited", required=True, help="the human-edited, saved deck")
    ap.add_argument("--threshold", "--threshold-inches", dest="threshold", type=float, default=0.03,
                    help="ignore position/size changes below this many inches (default 0.03)")
    ap.add_argument("--pos-tol", type=float, default=0.30,
                    help="proximity window in inches for text-edit recovery (default 0.30)")
    ap.add_argument("--slide", default=None, help="restrict to these slide numbers, e.g. 2,10,13")
    ap.add_argument("--json", action="store_true", help="emit machine-readable JSON")
    ap.add_argument("--hide-low-confidence", action="store_true",
                    help="drop low-confidence rows (equation/decorative re-serialization, title autofit)")
    a = ap.parse_args()

    if a.threshold < 0 or a.pos_tol < 0:
        print("error: --threshold and --pos-tol must be non-negative", file=sys.stderr)
        sys.exit(2)
    only = _parse_slides(a.slide)

    ref = _load(a.reference, "reference")
    ed = _load(a.edited, "edited")
    th, tol = a.threshold * EMU, a.pos_tol * EMU

    nref, ned = len(ref.slides), len(ed.slides)
    # Positional assumption: extra/missing slides are the trailing ones.
    added_slides = list(range(nref + 1, ned + 1)) if ned > nref else []
    deleted_slides = list(range(ned + 1, nref + 1)) if nref > ned else []
    slide_count = {"reference": nref, "edited": ned,
                   "added_slides": added_slides, "deleted_slides": deleted_slides,
                   "note": "extra/missing slides assumed to be trailing; verify if slides were inserted mid-deck"}

    changes = []
    for i, (rs, es) in enumerate(zip(ref.slides, ed.slides), 1):
        if only and i not in only:
            continue
        changes.extend(diff_slide(i, rs.shapes, es.shapes, th, tol))
    for s in added_slides:
        if not only or s in only:
            changes.append(_record(s, "slide_added", "unmatched", "high",
                                   "slide present in edited but not in reference"))
    for s in deleted_slides:
        if not only or s in only:
            changes.append(_record(s, "slide_deleted", "unmatched", "high",
                                   "slide present in reference but not in edited"))

    if a.hide_low_confidence:
        changes = [c for c in changes if c["confidence"] != "low"]

    if a.json:
        print(json.dumps({"slide_count": slide_count, "changes": changes},
                         ensure_ascii=False, indent=2))
        return

    # text output
    if slide_count["added_slides"] or slide_count["deleted_slides"]:
        print(f"SLIDE COUNT: reference={nref} edited={ned}  "
              f"added={slide_count['added_slides'] or '-'}  deleted={slide_count['deleted_slides'] or '-'}")
        print("  (extra/missing slides assumed trailing; verify if inserted mid-deck)")
    by_slide = defaultdict(list)
    for c in changes:
        by_slide[c["slide"]].append(c)
    if not by_slide:
        print("No differences above threshold. Source reproduces the edited deck.")
        return
    C = {"high": "HIGH", "medium": "med ", "low": "low "}
    for s in sorted(by_slide):
        print(f"SLIDE {s}:")
        for c in by_slide[s]:
            tag = "  [suspected noise]" if c["suspected_noise"] else ""
            head = f"  [{C[c['confidence']]}] {c['kind']:16} via {c['matched_by']:9}"
            bits = []
            if "text_from" in c:
                bits.append(f'text "{c["text_from"][:44]}" -> "{c["text_to"][:44]}"')
            for d in c.get("deltas", []):
                bits.append(f'{d["prop"]} {d["from"]}->{d["to"]}')
            if c["kind"] == "added_in_edited" and "after_bbox" in c:
                b = c["after_bbox"]; bits.append(f'@({b["x"]},{b["y"]}) {b["w"]}x{b["h"]}')
            if c["kind"] == "removed_in_edited" and "before_bbox" in c:
                b = c["before_bbox"]; bits.append(f'was @({b["x"]},{b["y"]}) {b["w"]}x{b["h"]}')
            print(head + "  " + "; ".join(bits) + tag)
            print(f"          reason: {c['reason']}")
    total = len(changes)
    print(f"=== {total} change(s); units=inches, threshold={a.threshold}in. "
          "Filter with --hide-low-confidence; low = equation/decorative re-serialization or title autofit. ===")


if __name__ == "__main__":
    main()
