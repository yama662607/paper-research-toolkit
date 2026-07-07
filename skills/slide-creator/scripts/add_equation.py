#!/usr/bin/env -S uv run --quiet
# /// script
# requires-python = ">=3.11"
# dependencies = ["python-pptx>=0.6.23", "latex2mathml>=3.77", "lxml>=5.0"]
# ///
"""Inject a native (editable) math equation into a pptx slide.

Pipeline: LaTeX -> MathML (latex2mathml) -> OMML (mathml2omml via bun)
-> wrap in a14:m -> append to a new textbox paragraph.

Existing OMML can also be injected directly with --omml/--omml-file. This is
mainly for round-trip sync, where PowerPoint-edited equations are preserved as
OMML instead of guessed back into LaTeX.

PowerPoint's DrawingML paragraphs don't accept m:oMath directly; the Office
2010 a14:m wrapper is required. We deliberately omit mc:AlternateContent
(PowerPoint accepts the bare form, and python-pptx mis-handles shapes inside
AlternateContent on later reads).
"""
import argparse
import re
import subprocess
import sys
from pathlib import Path

from lxml import etree

SCRIPT_DIR = Path(__file__).resolve().parent
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_A14 = "http://schemas.microsoft.com/office/drawing/2010/main"
NS_M = "http://schemas.openxmlformats.org/officeDocument/2006/math"
ENTITY_RE = re.compile(r"&(?!#\d+;|#x[0-9A-Fa-f]+;|amp;|lt;|gt;|quot;|apos;)")


def escape_omml_text_nodes(omml: str) -> str:
    """Escape raw text emitted inside <m:t> nodes by some MathML converters.

    mathml2omml can emit literal "<" in m:t text for inequalities. That makes
    the OMML fragment invalid XML even though the intended equation is valid.
    Limit the fix to text-node contents so real OMML tags are untouched.
    """

    def repl(match: re.Match[str]) -> str:
        attrs, body = match.group(1), match.group(2)
        body = ENTITY_RE.sub("&amp;", body)
        body = body.replace("<", "&lt;").replace(">", "&gt;")
        return f"<m:t{attrs}>{body}</m:t>"

    return re.sub(r"<m:t((?:\s[^>]*)?)>(.*?)</m:t>", repl, omml, flags=re.DOTALL)


def latex_to_omml(latex: str, display: bool) -> str:
    import latex2mathml.converter

    mathml = latex2mathml.converter.convert(
        latex, display="block" if display else "inline"
    )
    node = SCRIPT_DIR / "omml" / "convert.mjs"
    if not (SCRIPT_DIR / "omml" / "node_modules").exists():
        sys.exit(
            "error: equation converter not set up. Run once:\n"
            f"  cd {SCRIPT_DIR / 'omml'} && bun install"
        )
    proc = subprocess.run(
        ["bun", str(node)], input=mathml, capture_output=True, text=True
    )
    if proc.returncode != 0:
        sys.exit(f"error: MathML->OMML conversion failed:\n{proc.stderr}")
    return escape_omml_text_nodes(proc.stdout.strip())


def build_a14_wrapper(omml: str, display: bool) -> etree._Element:
    """Wrap OMML in <a14:m>, with m:oMathPara for display equations."""
    omml = omml.strip()
    try:
        root = etree.fromstring(omml.encode())
        if etree.QName(root).localname == "m" and root.nsmap.get("a14") == NS_A14:
            return root
    except etree.XMLSyntaxError:
        pass

    inner = omml if not display or "<m:oMathPara" in omml else (
        f'<m:oMathPara xmlns:m="{NS_M}">{omml}</m:oMathPara>'
    )
    xml = (
        f'<a14:m xmlns:a14="{NS_A14}" xmlns:m="{NS_M}" xmlns:a="{NS_A}">'
        f"{inner}</a14:m>"
    )
    try:
        return etree.fromstring(xml.encode())
    except etree.XMLSyntaxError as e:
        sys.exit(f"error: generated OMML is not well-formed XML: {e}\n---\n{xml}")


def main() -> None:
    p = argparse.ArgumentParser(description=__doc__)
    p.add_argument("pptx", help="target .pptx (modified in place)")
    p.add_argument("--slide", type=int, required=True, help="slide number (1-based)")
    src = p.add_mutually_exclusive_group(required=True)
    src.add_argument("--latex", help="LaTeX source of the equation")
    src.add_argument("--latex-file", help="file containing LaTeX source")
    src.add_argument("--omml", help="existing OMML fragment to inject")
    src.add_argument("--omml-file", help="file containing an existing OMML fragment")
    p.add_argument("--x", type=float, default=1.0, help="left, inches")
    p.add_argument("--y", type=float, default=2.0, help="top, inches")
    p.add_argument("--w", type=float, default=6.0, help="width, inches")
    p.add_argument("--h", type=float, default=1.2, help="height, inches")
    p.add_argument("--font-size", type=int, default=24, help="points")
    p.add_argument("--source-id",
                   help="stable slide-creator id; writes object name scid:<id>")
    p.add_argument("--object-name",
                   help="explicit PowerPoint object name for the equation box")
    p.add_argument("--inline", action="store_true",
                   help="inline equation (no oMathPara block wrapper)")
    p.add_argument("--emit-xml", action="store_true",
                   help="print the a14:m fragment and exit without modifying the file")
    args = p.parse_args()

    display = not args.inline
    if args.omml or args.omml_file:
        omml = args.omml if args.omml else Path(args.omml_file).read_text().strip()
    else:
        latex = args.latex if args.latex else Path(args.latex_file).read_text().strip()
        omml = latex_to_omml(latex, display)
    fragment = build_a14_wrapper(omml, display)

    if args.emit_xml:
        print(etree.tostring(fragment, pretty_print=True).decode())
        return

    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation(args.pptx)
    if not (1 <= args.slide <= len(prs.slides)):
        sys.exit(f"error: slide {args.slide} out of range (deck has {len(prs.slides)})")
    slide = prs.slides[args.slide - 1]

    box = slide.shapes.add_textbox(
        Inches(args.x), Inches(args.y), Inches(args.w), Inches(args.h)
    )
    object_name = args.object_name
    if args.source_id:
        object_name = object_name or f"scid:{args.source_id}"
    if object_name:
        cnv = box._element.xpath(".//p:cNvPr")[0]
        cnv.set("name", object_name)
        cnv.set("descr", object_name)
        cnv.set("title", object_name)

    para = box.text_frame.paragraphs[0]
    # Font size on the paragraph's endParaRPr scales the equation display.
    para.font.size = Pt(args.font_size)
    para._p.append(fragment)

    prs.save(args.pptx)
    print(f"ok: equation injected on slide {args.slide} of {args.pptx}")
    print("verify: uv run scripts/verify_deck.py + open in PowerPoint "
          "(equation must be double-click editable)")


if __name__ == "__main__":
    main()
