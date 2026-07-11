#!/usr/bin/env -S uv run --quiet
# /// script
# requires-python = ">=3.11"
# dependencies = ["python-pptx>=0.6.23", "lxml>=5.0"]
# ///
"""Structural verification of a pptx. The only check that is safe at every
pipeline stage (it never renders, so it cannot trip LibreOffice's
video/animation bugs).

Checks: ZIP integrity, package/content-type consistency, python-pptx re-open,
relationship targets and references, media health via ffprobe, duplicate shape
ids, animation target ids, placeholder debris in text, fake list markers.
With --min-font-size, also enforces an explicit audience-readable text floor;
--text dumps per-slide text instead.
"""
import argparse
from collections import Counter
import json
import math
import re
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path

PLACEHOLDER_RE = re.compile(
    r"\b(lorem|ipsum|TODO|FIXME|XXX{1,}|\[insert|placeholder|ここに(入力|テキスト))",
    re.IGNORECASE,
)
FAKE_LIST_MARKER_RE = re.compile(
    r"^\s*(?:"
    r"[•‣⁃◦▪▫●・]\s*"
    r"|[-–—*]\s+"
    r"|(?:[0-9０-９]{1,2}|[A-Za-zＡ-Ｚａ-ｚ])[\.\)．）]\s+"
    r"|[①②③④⑤⑥⑦⑧⑨⑩⑪⑫⑬⑭⑮⑯⑰⑱⑲⑳]\s*"
    r")"
)
BULLET_TAGS = ("buChar", "buAutoNum", "buBlip")
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def fail(msg: str, problems: list[str]) -> None:
    problems.append(msg)
    print(f"  FAIL {msg}")


def ok(msg: str) -> None:
    print(f"  ok   {msg}")


def _has_native_list_marker(el, ns: dict[str, str]) -> bool:
    return any(el.find(f"a:{tag}", namespaces=ns) is not None
               for tag in BULLET_TAGS)


def _paragraph_has_native_list_marker(para, text_body, ns: dict[str, str]) -> bool:
    ppr = para.find("a:pPr", namespaces=ns)
    level = 0
    if ppr is not None:
        if _has_native_list_marker(ppr, ns):
            return True
        if ppr.find("a:buNone", namespaces=ns) is not None:
            return False
        raw_level = ppr.get("lvl")
        if raw_level and raw_level.isdigit():
            level = int(raw_level)

    # PowerPoint can inherit bullets from the text body's list style.
    lst_style = text_body.find("a:lstStyle", namespaces=ns)
    if lst_style is None:
        return False
    lvl_ppr = lst_style.find(f"a:lvl{level + 1}pPr", namespaces=ns)
    if lvl_ppr is None:
        return False
    if lvl_ppr.find("a:buNone", namespaces=ns) is not None:
        return False
    return _has_native_list_marker(lvl_ppr, ns)


def _rels_name_for_part(part_name: str) -> str:
    part = Path(part_name)
    return (part.parent / "_rels" / f"{part.name}.rels").as_posix()


def _content_type_extension(part_name: str) -> str:
    if part_name.endswith(".rels"):
        return "rels"
    return Path(part_name).suffix.lower().lstrip(".")


def _compile_patterns(raw_patterns: list[str], parser: argparse.ArgumentParser,
                      option: str) -> list[re.Pattern[str]]:
    compiled = []
    for raw in raw_patterns:
        try:
            compiled.append(re.compile(raw, re.IGNORECASE))
        except re.error as exc:
            parser.error(f"{option} contains invalid regex {raw!r}: {exc}")
    return compiled


def _matches_any(patterns: list[re.Pattern[str]], value: str) -> bool:
    return any(pattern.search(value) for pattern in patterns)


def _iter_text_frames(shapes, prefix: str = ""):
    """Yield (PowerPoint name, diagnostic label, text frame)."""
    for shape in shapes:
        shape_name = shape.name or "<unnamed>"
        label = f"{prefix}{shape_name}"
        if shape.has_text_frame:
            yield shape_name, label, shape.text_frame
        if shape.has_table:
            for row_idx, row in enumerate(shape.table.rows, 1):
                for col_idx, cell in enumerate(row.cells, 1):
                    yield (
                        shape_name,
                        f"{label} cell {row_idx},{col_idx}",
                        cell.text_frame,
                    )
        child_shapes = getattr(shape, "shapes", None)
        if child_shapes is not None:
            yield from _iter_text_frames(child_shapes, prefix=f"{label} / ")


def _paragraph_segments(para) -> list[tuple[str, float | None]]:
    """Return visible a:r/a:fld text and explicit size in points.

    python-pptx Paragraph.runs omits PowerPoint fields such as slide numbers
    and dates, so inspect the paragraph XML directly.
    """
    from pptx.oxml.ns import qn

    segments = []
    for node in para._p.xpath("./a:r | ./a:fld"):
        text = "".join(t.text or "" for t in node.xpath(".//a:t"))
        if not text.strip():
            continue
        rpr = node.find(qn("a:rPr"))
        raw_size = rpr.get("sz") if rpr is not None else None
        try:
            size_pt = float(raw_size) / 100 if raw_size is not None else None
        except (TypeError, ValueError):
            size_pt = None
        if size_pt is not None and (not math.isfinite(size_pt) or size_pt <= 0):
            size_pt = None
        segments.append((text, size_pt))
    return segments


def main() -> None:
    p = argparse.ArgumentParser(description=__doc__)
    p.add_argument("pptx")
    p.add_argument("--text", action="store_true",
                   help="dump per-slide text and exit")
    p.add_argument(
        "--min-font-size",
        type=float,
        metavar="PT",
        help=("fail non-exempt text runs below PT; text with no explicit run "
              "size also fails because the floor cannot be proven"),
    )
    p.add_argument(
        "--allow-small-font-regex",
        action="append",
        default=[],
        metavar="REGEX",
        help=("exempt a paragraph whose text matches REGEX; repeat for narrow "
              "citation/page-number exceptions"),
    )
    p.add_argument(
        "--allow-small-font-shape-regex",
        action="append",
        default=[],
        metavar="REGEX",
        help=("exempt text in a shape whose PowerPoint name matches REGEX; "
              "repeat as needed"),
    )
    args = p.parse_args()
    if (args.min_font_size is not None
            and (not math.isfinite(args.min_font_size)
                 or args.min_font_size <= 0)):
        p.error("--min-font-size must be a finite number greater than zero")
    text_exemptions = _compile_patterns(
        args.allow_small_font_regex, p, "--allow-small-font-regex"
    )
    shape_exemptions = _compile_patterns(
        args.allow_small_font_shape_regex, p,
        "--allow-small-font-shape-regex"
    )
    path = Path(args.pptx)
    problems: list[str] = []

    # --- ZIP layer -------------------------------------------------------
    if not zipfile.is_zipfile(path):
        sys.exit(f"FAIL: {path} is not a ZIP archive")
    zf = zipfile.ZipFile(path)
    if bad := zf.testzip():
        sys.exit(f"FAIL: corrupt ZIP member: {bad}")
    entries = zf.namelist()
    names = set(entries)
    file_names = {n for n in names if not n.endswith("/")}
    if "[Content_Types].xml" not in names:
        sys.exit("FAIL: [Content_Types].xml missing at archive root "
                 "(usually a re-zip that included the parent directory)")

    from lxml import etree
    entry_counts = Counter(entries)
    duplicate_entries = sorted(
        name for name, count in entry_counts.items()
        if count > 1 and not name.endswith("/")
    )
    try:
        ct_root = etree.fromstring(zf.read("[Content_Types].xml"))
    except etree.XMLSyntaxError as e:
        sys.exit(f"FAIL: [Content_Types].xml is not well-formed XML: {e}")
    CT_NS = {"ct": "http://schemas.openxmlformats.org/package/2006/content-types"}
    overrides_by_part: dict[str, list[str | None]] = {}
    for override in ct_root.findall("ct:Override", namespaces=CT_NS):
        part_name = override.get("PartName", "").lstrip("/")
        if part_name:
            overrides_by_part.setdefault(part_name, []).append(
                override.get("ContentType")
            )
    missing_overrides = [
        part_name for part_name in overrides_by_part
        if part_name not in file_names
    ]
    duplicate_overrides = sorted(
        part_name for part_name, content_types in overrides_by_part.items()
        if len(content_types) > 1
    )
    defaults_by_ext: dict[str, list[str | None]] = {}
    for default in ct_root.findall("ct:Default", namespaces=CT_NS):
        ext = (default.get("Extension") or "").lower()
        if ext:
            defaults_by_ext.setdefault(ext, []).append(default.get("ContentType"))
    duplicate_defaults = sorted(
        ext for ext, content_types in defaults_by_ext.items()
        if len(content_types) > 1
    )
    parts_without_content_type = []
    for part_name in sorted(file_names):
        if part_name == "[Content_Types].xml":
            continue
        ext = _content_type_extension(part_name)
        if part_name not in overrides_by_part and ext not in defaults_by_ext:
            parts_without_content_type.append(part_name)

    # --- python-pptx re-open ---------------------------------------------
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        n_slides = len(prs.slides)
    except Exception as e:
        sys.exit(f"FAIL: python-pptx cannot open the file: {e}")

    if args.text:
        for i, slide in enumerate(prs.slides, 1):
            print(f"## Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(r.text for r in para.runs).strip()
                        if t:
                            print(t)
            print()
        return

    print(f"verifying {path} ({n_slides} slides)")
    ok("ZIP integrity, [Content_Types].xml at root, python-pptx opens")
    if duplicate_entries:
        fail(f"duplicate ZIP entries: {duplicate_entries}", problems)
    else:
        ok("ZIP entries are unique")
    if missing_overrides:
        fail("[Content_Types].xml overrides point to missing package parts: "
             f"{missing_overrides}", problems)
    else:
        ok("[Content_Types].xml overrides point to real package parts")
    if duplicate_overrides:
        fail(f"duplicate [Content_Types].xml overrides: {duplicate_overrides}",
             problems)
    else:
        ok("[Content_Types].xml overrides are unique")
    if duplicate_defaults:
        fail(f"duplicate [Content_Types].xml defaults: {duplicate_defaults}",
             problems)
    else:
        ok("[Content_Types].xml defaults are unique")
    if parts_without_content_type:
        fail("package parts without content type default/override: "
             f"{parts_without_content_type}", problems)
    else:
        ok("all package parts have a content type")

    # --- relationship targets exist --------------------------------------
    NS = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "r": R_NS,
    }
    missing = []
    duplicate_rel_ids = []
    bad_external_targets = []
    rel_ids_by_rels: dict[str, set[str]] = {}
    for rels_name in [n for n in names if n.endswith(".rels")]:
        base = Path(rels_name).parent.parent  # e.g. ppt/slides
        root = etree.fromstring(zf.read(rels_name))
        ids = [rel.get("Id") for rel in root if rel.get("Id")]
        rel_ids_by_rels[rels_name] = set(ids)
        dups = sorted({rid for rid in ids if ids.count(rid) > 1})
        if dups:
            duplicate_rel_ids.append(f"{rels_name}: {dups}")
        for rel in root:
            target, mode = rel.get("Target", ""), rel.get("TargetMode", "Internal")
            if not target.strip():
                if mode == "External":
                    bad_external_targets.append(
                        f"{rels_name}: {rel.get('Id', '<no id>')} has empty external Target"
                    )
                else:
                    missing.append(
                        f"{rels_name}: {rel.get('Id', '<no id>')} has empty Target"
                    )
                continue
            if mode == "External":
                continue
            resolved = (base / target).as_posix()
            parts = []
            for seg in resolved.split("/"):
                if seg == "..":
                    if parts:
                        parts.pop()
                elif seg not in (".", ""):
                    parts.append(seg)
            if "/".join(parts) not in names:
                missing.append(f"{rels_name}: {target}")
    if missing:
        fail(f"unresolved relationship targets: {missing}", problems)
    else:
        ok("all internal relationship targets resolve")
    if duplicate_rel_ids:
        fail(f"duplicate relationship ids: {duplicate_rel_ids}", problems)
    else:
        ok("relationship ids unique within each .rels part")
    if bad_external_targets:
        fail(f"empty external relationship targets: {bad_external_targets}",
             problems)
    else:
        ok("external relationship targets are non-empty")

    bad_rel_refs = []
    xml_parts = [
        n for n in names
        if n.endswith(".xml")
        and not n.endswith(".rels")
        and n != "[Content_Types].xml"
    ]
    for part_name in xml_parts:
        rels_name = _rels_name_for_part(part_name)
        rel_ids = rel_ids_by_rels.get(rels_name)
        try:
            root = etree.fromstring(zf.read(part_name))
        except etree.XMLSyntaxError:
            continue
        for el in root.iter():
            for attr_name, value in el.attrib.items():
                attr = etree.QName(attr_name)
                if attr.namespace != R_NS:
                    continue
                where = (
                    f"{part_name}: {attr.localname} on "
                    f"<{etree.QName(el).localname}>"
                )
                if not value:
                    bad_rel_refs.append(f"{where} is empty")
                elif rel_ids is None:
                    bad_rel_refs.append(f"{where}={value!r} but {rels_name} is missing")
                elif value not in rel_ids:
                    bad_rel_refs.append(f"{where}={value!r} is not defined in {rels_name}")
    if bad_rel_refs:
        fail(f"unresolved relationship references in XML: {bad_rel_refs}", problems)
    else:
        ok("all relationship references in XML resolve")

    # --- shape ids and animation targets ----------------------------------
    dup_report = []
    bad_animation_targets = []
    pres_root = etree.fromstring(zf.read("ppt/presentation.xml"))
    prels_root = etree.fromstring(zf.read("ppt/_rels/presentation.xml.rels"))
    rid_to_target = {rel.get("Id"): rel.get("Target") for rel in prels_root}
    slide_names = []
    sldlst = pres_root.find("p:sldIdLst", namespaces=NS)
    if sldlst is not None:
        for sld in sldlst:
            rid = sld.get(f"{{{NS['r']}}}id")
            target = rid_to_target.get(rid)
            if target:
                slide_names.append((Path("ppt") / target).as_posix())
    for i, slide in enumerate(prs.slides, 1):
        ids = [s.shape_id for s in slide.shapes]
        dups = {x for x in ids if ids.count(x) > 1}
        if dups:
            dup_report.append(f"slide {i}: {sorted(dups)}")
        if i <= len(slide_names):
            root = etree.fromstring(zf.read(slide_names[i - 1]))
            xml_ids = {
                sp.get("id")
                for sp in root.findall(".//p:cNvPr", namespaces=NS)
                if sp.get("id")
            }
            for tgt in root.findall(".//p:timing//p:spTgt", namespaces=NS):
                spid = tgt.get("spid")
                if spid and spid not in xml_ids:
                    bad_animation_targets.append(
                        f"slide {i}: spid {spid} not in shape ids")
    if dup_report:
        fail(f"duplicate shape ids (animations may target wrong shapes): "
             f"{dup_report}", problems)
    else:
        ok("shape ids unique per slide")
    if bad_animation_targets:
        fail(f"animation targets missing shapes: {bad_animation_targets}",
             problems)
    else:
        ok("animation targets resolve to slide shapes")

    # --- placeholder debris ------------------------------------------------
    debris = []
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs)
                    if PLACEHOLDER_RE.search(t):
                        debris.append(f"slide {i}: {t.strip()[:60]!r}")
    if debris:
        fail(f"placeholder-looking text: {debris}", problems)
    else:
        ok("no placeholder debris in text")

    # --- native-list enforcement -------------------------------------------
    fake_list_markers = []
    for i, slide_name in enumerate(slide_names, 1):
        root = etree.fromstring(zf.read(slide_name))
        for tx_body in root.xpath(".//*[local-name()='txBody']"):
            for para in tx_body.findall("a:p", namespaces=NS):
                text = "".join(
                    node.text or ""
                    for node in para.findall(".//a:t", namespaces=NS)
                )
                if not text.strip() or not FAKE_LIST_MARKER_RE.match(text):
                    continue
                if _paragraph_has_native_list_marker(para, tx_body, NS):
                    continue
                fake_list_markers.append(f"slide {i}: {text.strip()[:60]!r}")
    if fake_list_markers:
        fail("manual list-marker-looking text; use native PowerPoint bullets/"
             f"numbering instead: {fake_list_markers}", problems)
    else:
        ok("no manual list markers masquerading as bullets/numbering")

    # --- audience-readable font floor --------------------------------------
    if args.min_font_size is not None:
        undersized = []
        unknown_size = []
        for i, slide in enumerate(prs.slides, 1):
            for shape_name, text_label, text_frame in _iter_text_frames(slide.shapes):
                if _matches_any(shape_exemptions, shape_name):
                    continue
                for para in text_frame.paragraphs:
                    para_text = para.text.strip()
                    if not para_text or _matches_any(text_exemptions, para_text):
                        continue
                    segments = _paragraph_segments(para)
                    if not segments:
                        unknown_size.append(
                            f"slide {i}, shape {text_label!r}: {para_text[:60]!r}"
                        )
                        continue
                    for segment_text, size_pt in segments:
                        snippet = segment_text.strip()[:60]
                        where = f"slide {i}, shape {text_label!r}: {snippet!r}"
                        if size_pt is None:
                            unknown_size.append(where)
                        elif size_pt + 1e-6 < args.min_font_size:
                            undersized.append(f"{where} ({size_pt:g} pt)")
        if undersized:
            fail(
                f"text below {args.min_font_size:g} pt audience floor: "
                f"{undersized}", problems
            )
        else:
            ok(f"no non-exempt text below {args.min_font_size:g} pt")
        if unknown_size:
            fail(
                "text has no explicit run font size, so the audience floor "
                f"cannot be proven: {unknown_size}", problems
            )
        else:
            ok("all non-exempt text has an explicit run font size")

    # --- media health -------------------------------------------------------
    media = [n for n in names if n.startswith("ppt/media/")
             and n.lower().endswith((".mp4", ".m4v", ".mov", ".avi", ".wmv"))]
    if media:
        if shutil.which("ffprobe") is None:
            fail("videos present but ffprobe unavailable — cannot verify", problems)
        else:
            import tempfile
            for m in media:
                with tempfile.NamedTemporaryFile(suffix=Path(m).suffix) as tmp:
                    tmp.write(zf.read(m))
                    tmp.flush()
                    proc = subprocess.run(
                        ["ffprobe", "-v", "error", "-show_streams", "-of", "json",
                         tmp.name], capture_output=True, text=True)
                    if proc.returncode != 0:
                        fail(f"{m}: ffprobe error: {proc.stderr.strip()[:120]}",
                             problems)
                        continue
                    v = next((s for s in json.loads(proc.stdout)["streams"]
                              if s["codec_type"] == "video"), None)
                    codec = v.get("codec_name") if v else "none"
                    pix = v.get("pix_fmt") if v else "-"
                    if codec != "h264" or pix != "yuv420p":
                        fail(f"{m}: codec={codec}/{pix} (PowerPoint wants "
                             f"h264/yuv420p — re-embed with --normalize)", problems)
                    else:
                        ok(f"{m}: h264/yuv420p")
    else:
        ok("no embedded videos (nothing to probe)")

    # --- summary --------------------------------------------------------------
    print()
    if problems:
        print(f"RESULT: {len(problems)} problem(s) — fix before delivering")
        sys.exit(1)
    print("RESULT: all checks passed")


if __name__ == "__main__":
    main()
