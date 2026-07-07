#!/usr/bin/env -S uv run --quiet
# /// script
# requires-python = ">=3.11"
# dependencies = ["python-pptx>=0.6.23", "lxml>=5.0"]
# ///
"""Slide transitions and shape animations via verified XML templates.

Shape-animation XML is fragile: PowerPoint silently drops malformed timing
trees. This tool therefore only instantiates proven templates
(assets/timing-templates/) with substituted ids/durations, and refuses to
merge into timing trees it does not understand.

Subcommands:
  shapes      list shape ids/names of a slide (targets for `effect`)
  transition  set the slide transition (fade | wipe | push | morph)
  effect      add an entrance animation to a shape
"""
import argparse
import string
import sys
from pathlib import Path

from lxml import etree

SCRIPT_DIR = Path(__file__).resolve().parent
TEMPLATE_DIR = SCRIPT_DIR.parent / "assets" / "timing-templates"

P = "http://schemas.openxmlformats.org/presentationml/2006/main"
MC = "http://schemas.openxmlformats.org/markup-compatibility/2006"
P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
P159 = "http://schemas.microsoft.com/office/powerpoint/2015/09/main"


def q(tag: str) -> str:
    return f"{{{P}}}{tag}"


# ---------------------------------------------------------------- helpers

def load_prs(path: str):
    from pptx import Presentation
    return Presentation(path)


def get_slide(prs, num: int):
    if not (1 <= num <= len(prs.slides)):
        sys.exit(f"error: slide {num} out of range (deck has {len(prs.slides)})")
    return prs.slides[num - 1]


def sld_order_index(sld, tag_local: str) -> int:
    """Insertion index that keeps p:sld child order: cSld, clrMapOvr,
    transition, timing, extLst."""
    order = ["cSld", "clrMapOvr", "transition", "timing", "extLst"]
    rank = order.index(tag_local)
    idx = 0
    for child in sld:
        local = etree.QName(child).localname
        # AlternateContent wrapping a transition sorts as 'transition'
        if local == "AlternateContent":
            inner = child.find(f".//{q('transition')}")
            local = "transition" if inner is not None else local
        if local in order and order.index(local) <= rank:
            idx = list(sld).index(child) + 1
    return idx


# ------------------------------------------------------------- transition

PLAIN_TRANSITIONS = {
    "fade": f'<p:transition xmlns:p="{P}" spd="med"><p:fade/></p:transition>',
    "wipe": f'<p:transition xmlns:p="{P}" spd="med"><p:wipe/></p:transition>',
}


def ext_transition(kind: str, duration_ms: int) -> str:
    """p14/p159 transitions need mc:AlternateContent with a fade fallback,
    or old PowerPoint / LibreOffice shows a repair dialog."""
    if kind == "push":
        choice_ns, requires, body = P14, "p14", "<p14:push/>"
    elif kind == "morph":
        choice_ns, requires, body = P159, "p159", '<p159:morph option="byObject"/>'
    else:
        raise ValueError(kind)
    return (
        f'<mc:AlternateContent xmlns:mc="{MC}">'
        f'<mc:Choice xmlns:{requires}="{choice_ns}" Requires="{requires}">'
        f'<p:transition xmlns:p="{P}" xmlns:p14="{P14}" spd="med" p14:dur="{duration_ms}">'
        f"{body}</p:transition></mc:Choice>"
        f'<mc:Fallback><p:transition xmlns:p="{P}" spd="med"><p:fade/></p:transition>'
        f"</mc:Fallback></mc:AlternateContent>"
    )


def remove_existing_transition(sld) -> None:
    for child in list(sld):
        local = etree.QName(child).localname
        if local == "transition":
            sld.remove(child)
        elif local == "AlternateContent" and child.find(f".//{q('transition')}") is not None:
            sld.remove(child)


def cmd_transition(args) -> None:
    prs = load_prs(args.pptx)
    targets = (
        list(range(1, len(prs.slides) + 1)) if args.all else [args.slide]
    )
    if targets == [None]:
        sys.exit("error: give --slide N or --all")
    if args.type in PLAIN_TRANSITIONS:
        xml = PLAIN_TRANSITIONS[args.type]
    else:
        xml = ext_transition(args.type, args.duration)
    for num in targets:
        slide = get_slide(prs, num)
        sld = slide._element
        remove_existing_transition(sld)
        frag = etree.fromstring(xml.encode())
        sld.insert(sld_order_index(sld, "transition"), frag)
    prs.save(args.pptx)
    label = "all slides" if args.all else f"slide {args.slide}"
    print(f"ok: transition '{args.type}' set on {label}")
    if args.type == "morph":
        print("note: morph matches shapes BY NAME across consecutive slides — "
              "ensure moving shapes share p:cNvPr name; remove empty placeholders.")


# ----------------------------------------------------------------- shapes

def cmd_shapes(args) -> None:
    prs = load_prs(args.pptx)
    slide = get_slide(prs, args.slide)
    print(f"slide {args.slide} shapes (spid  name  type):")
    for shape in slide.shapes:
        print(f"  id={shape.shape_id:<4} name={shape.name!r}  {shape.shape_type}")


# ----------------------------------------------------------------- effect

ROOT_SCAFFOLD = f"""<p:timing xmlns:p="{P}">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst/>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>"""

MAIN_SEQ = f"""<p:seq xmlns:p="{P}" concurrent="1" nextAc="seek">
  <p:cTn id="${{ID}}" dur="indefinite" nodeType="mainSeq">
    <p:childTnLst/>
  </p:cTn>
  <p:prevCondLst>
    <p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
  </p:prevCondLst>
  <p:nextCondLst>
    <p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
  </p:nextCondLst>
</p:seq>"""

CLICK_GROUP = f"""<p:par xmlns:p="{P}">
  <p:cTn id="${{ID1}}" fill="hold">
    <p:stCondLst><p:cond delay="${{GROUP_DELAY}}"/></p:stCondLst>
    <p:childTnLst>
      <p:par>
        <p:cTn id="${{ID2}}" fill="hold">
          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
          <p:childTnLst/>
        </p:cTn>
      </p:par>
    </p:childTnLst>
  </p:cTn>
</p:par>"""


def max_ctn_id(timing) -> int:
    ids = [int(c.get("id")) for c in timing.iter(q("cTn"))
           if c.get("id", "").isdigit()]
    return max(ids, default=0)


def next_ids(timing, n: int) -> list[str]:
    base = max_ctn_id(timing)
    return [str(base + i + 1) for i in range(n)]


def cmd_effect(args) -> None:
    template_path = TEMPLATE_DIR / f"effect-{args.effect}.xml"
    if not template_path.exists():
        avail = ", ".join(p.stem.removeprefix("effect-")
                          for p in TEMPLATE_DIR.glob("effect-*.xml"))
        sys.exit(f"error: unknown effect '{args.effect}' (available: {avail})")

    prs = load_prs(args.pptx)
    slide = get_slide(prs, args.slide)
    sld = slide._element

    ids = {s.shape_id for s in slide.shapes}
    if args.spid not in ids:
        sys.exit(f"error: no shape with id {args.spid} on slide {args.slide} "
                 f"(run: animate.py shapes {args.pptx} --slide {args.slide})")

    # refuse AlternateContent-wrapped timing (unknown territory)
    for ac in sld.findall(f"{{{MC}}}AlternateContent"):
        if ac.find(f".//{q('timing')}") is not None:
            sys.exit("error: slide has an mc:AlternateContent-wrapped timing tree "
                     "(non-standard animation). Refusing to merge — this is a known "
                     "corruption path. Rebuild this slide's animations from scratch "
                     "or leave it un-animated.")

    timing = sld.find(q("timing"))
    if timing is None:
        timing = etree.fromstring(ROOT_SCAFFOLD.encode())
        sld.insert(sld_order_index(sld, "timing"), timing)

    troot = None
    for ctn in timing.iter(q("cTn")):
        if ctn.get("nodeType") == "tmRoot":
            troot = ctn
            break
    if troot is None:
        sys.exit("error: existing <p:timing> has no tmRoot node — unrecognized "
                 "structure, refusing to modify.")
    troot_children = troot.find(q("childTnLst"))
    if troot_children is None:
        troot_children = etree.SubElement(troot, q("childTnLst"))

    # find or create the main sequence
    mainseq_ctn = None
    for ctn in timing.iter(q("cTn")):
        if ctn.get("nodeType") == "mainSeq":
            mainseq_ctn = ctn
            break
    if mainseq_ctn is None:
        (seq_id,) = next_ids(timing, 1)
        seq = etree.fromstring(
            string.Template(MAIN_SEQ).substitute(ID=seq_id).encode()
        )
        troot_children.append(seq)
        mainseq_ctn = seq.find(q("cTn"))
    mainseq_children = mainseq_ctn.find(q("childTnLst"))

    # decide placement per trigger
    trigger = args.trigger
    effect_delay = "0"
    if trigger in ("click", "auto"):
        id1, id2 = next_ids(timing, 2)
        group_delay = "indefinite" if trigger == "click" else "0"
        group = etree.fromstring(
            string.Template(CLICK_GROUP)
            .substitute(ID1=id1, ID2=id2, GROUP_DELAY=group_delay)
            .encode()
        )
        mainseq_children.append(group)
        inner = group.find(f".//{q('cTn')}[@id='{id2}']/{q('childTnLst')}")
    else:  # with / after — attach to the last existing group
        groups = mainseq_children.findall(q("par"))
        if not groups:
            sys.exit("error: --trigger with/after needs a preceding effect on this "
                     "slide; add one with --trigger click first.")
        inner_ctns = groups[-1].findall(f".//{q('cTn')}/{q('childTnLst')}/"
                                        f"{q('par')}/{q('cTn')}")
        inner = groups[-1].find(f"{q('cTn')}/{q('childTnLst')}/{q('par')}/"
                                f"{q('cTn')}/{q('childTnLst')}")
        if inner is None:
            sys.exit("error: last animation group has an unexpected shape; "
                     "refusing to modify.")
        if trigger == "after":
            effect_delay = str(args.delay)

    tid1, tid2, tid3 = next_ids(timing, 3)
    effect_xml = string.Template(template_path.read_text()).substitute(
        ID1=tid1, ID2=tid2, ID3=tid3,
        SPID=str(args.spid), DUR=str(args.duration),
        NODETYPE="clickEffect" if trigger == "click" else
                 ("afterEffect" if trigger == "after" else "withEffect"),
        EFFECT_DELAY=effect_delay,
    )
    inner.append(etree.fromstring(effect_xml.encode()))

    prs.save(args.pptx)
    print(f"ok: '{args.effect}' ({trigger}) on shape {args.spid}, "
          f"slide {args.slide} of {args.pptx}")
    print("verify structurally only (verify_deck.py); do NOT re-render via "
          "LibreOffice after animating. Confirm playback once in PowerPoint.")


# ------------------------------------------------------------------- main

def main() -> None:
    p = argparse.ArgumentParser(description=__doc__)
    sub = p.add_subparsers(dest="cmd", required=True)

    s = sub.add_parser("shapes", help="list shape ids on a slide")
    s.add_argument("pptx")
    s.add_argument("--slide", type=int, required=True)
    s.set_defaults(fn=cmd_shapes)

    t = sub.add_parser("transition", help="set slide transition")
    t.add_argument("pptx")
    t.add_argument("--slide", type=int)
    t.add_argument("--all", action="store_true")
    t.add_argument("--type", required=True,
                   choices=["fade", "wipe", "push", "morph"])
    t.add_argument("--duration", type=int, default=700,
                   help="ms (push/morph only)")
    t.set_defaults(fn=cmd_transition)

    e = sub.add_parser("effect", help="add entrance animation to a shape")
    e.add_argument("pptx")
    e.add_argument("--slide", type=int, required=True)
    e.add_argument("--spid", type=int, required=True,
                   help="shape id (see `shapes` subcommand)")
    e.add_argument("--effect", required=True,
                   help="appear | fade-in | wipe-in")
    e.add_argument("--trigger", default="click",
                   choices=["click", "with", "after", "auto"])
    e.add_argument("--delay", type=int, default=500,
                   help="ms (trigger=after only)")
    e.add_argument("--duration", type=int, default=500, help="ms")
    e.set_defaults(fn=cmd_effect)

    args = p.parse_args()
    args.fn(args)


if __name__ == "__main__":
    main()
