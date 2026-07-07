#!/usr/bin/env -S uv run --quiet
# /// script
# requires-python = ">=3.11"
# dependencies = ["python-pptx>=0.6.23"]
# ///
"""Embed a video into a pptx slide, safely.

Wraps python-pptx add_movie() and fills its gaps:
- optional ffmpeg normalization to H.264/AAC/yuv420p/CFR/faststart
- poster frame auto-extraction (first frame) unless --poster given
- explicit mime_type (add_movie defaults to video/unknown)
- --autoplay: rewrites the media timing condition after insertion
- refuses slides whose <p:timing> is wrapped in mc:AlternateContent
  (python-pptx would duplicate the timing tree and corrupt the file)
"""
import argparse
import json
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "mc": "http://schemas.openxmlformats.org/markup-compatibility/2006",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def run(cmd: list[str], **kw) -> subprocess.CompletedProcess:
    return subprocess.run(cmd, capture_output=True, text=True, **kw)


def ffprobe(path: str) -> dict:
    proc = run(["ffprobe", "-v", "error", "-show_format", "-show_streams",
                "-of", "json", path])
    if proc.returncode != 0:
        sys.exit(f"error: ffprobe failed on {path}:\n{proc.stderr}")
    return json.loads(proc.stdout)


def needs_normalization(info: dict) -> list[str]:
    reasons = []
    v = next((s for s in info["streams"] if s["codec_type"] == "video"), None)
    if v is None:
        sys.exit("error: no video stream found")
    if v.get("codec_name") != "h264":
        reasons.append(f"codec is {v.get('codec_name')} (need h264)")
    if v.get("pix_fmt") != "yuv420p":
        reasons.append(f"pix_fmt is {v.get('pix_fmt')} (need yuv420p)")
    if v.get("avg_frame_rate") != v.get("r_frame_rate"):
        reasons.append("variable frame rate suspected")
    a = next((s for s in info["streams"] if s["codec_type"] == "audio"), None)
    if a is not None and a.get("codec_name") != "aac":
        reasons.append(f"audio codec is {a.get('codec_name')} (need aac)")
    return reasons


def normalize(src: str, dst: str) -> None:
    proc = run([
        "ffmpeg", "-y", "-i", src,
        "-c:v", "libx264", "-profile:v", "high", "-preset", "medium", "-crf", "18",
        "-pix_fmt", "yuv420p", "-r", "30", "-fps_mode", "cfr",
        "-c:a", "aac", "-b:a", "192k", "-ar", "48000",
        "-movflags", "+faststart", dst,
    ])
    if proc.returncode != 0:
        sys.exit(f"error: ffmpeg normalization failed:\n{proc.stderr[-2000:]}")


def extract_poster(video: str, dst: str) -> None:
    proc = run(["ffmpeg", "-y", "-i", video, "-vf", "select=eq(n\\,0)",
                "-frames:v", "1", dst])
    if proc.returncode != 0:
        sys.exit(f"error: poster extraction failed:\n{proc.stderr[-2000:]}")


def remove_empty_media_hyperlinks(shape) -> int:
    """Remove add_movie()'s empty media action hyperlink, if present.

    python-pptx may emit <a:hlinkClick r:id="" action="ppaction://media"/>.
    PowerPoint often tolerates it, but it is not a resolvable relationship and
    makes structural validation weaker than it should be.
    """
    removed = 0
    rid_attr = f"{{{NS['r']}}}id"
    for hlink in list(shape._element.iter(f"{{{NS['a']}}}hlinkClick")):
        if hlink.get(rid_attr) == "" and hlink.get("action") == "ppaction://media":
            parent = hlink.getparent()
            if parent is not None:
                parent.remove(hlink)
                removed += 1
    return removed


def set_video_autoplay(slide_element, shape_id: int) -> None:
    timing = slide_element.find(f"{{{NS['p']}}}timing")
    if timing is None:
        sys.exit("error: add_movie did not create a timing node (unexpected)")

    shape_id_str = str(shape_id)
    media_nodes = []
    for node in timing.iter(f"{{{NS['p']}}}cMediaNode"):
        target = node.find(
            f"{{{NS['p']}}}tgtEl/{{{NS['p']}}}spTgt"
        )
        if target is not None and target.get("spid") == shape_id_str:
            media_nodes.append(node)
    if len(media_nodes) != 1:
        sys.exit(
            "error: expected exactly one media timing node for the inserted "
            f"video shape id {shape_id}, found {len(media_nodes)}"
        )

    conds = [
        cond for cond in media_nodes[0].iter(f"{{{NS['p']}}}cond")
        if cond.get("delay") == "indefinite"
    ]
    if len(conds) != 1:
        sys.exit(
            "error: expected exactly one click-to-play condition for the "
            f"inserted video shape id {shape_id}, found {len(conds)}"
        )

    cond = conds[0]
    cond.attrib.pop("delay", None)
    cond.set("evt", "onBegin")
    cond.set("delay", "0")


def main() -> None:
    p = argparse.ArgumentParser(description=__doc__)
    p.add_argument("pptx")
    p.add_argument("video")
    p.add_argument("--slide", type=int, required=True, help="1-based")
    p.add_argument("--x", type=float, default=1.0)
    p.add_argument("--y", type=float, default=1.0)
    p.add_argument("--w", type=float, default=5.0)
    p.add_argument("--h", type=float, default=2.8125)
    p.add_argument("--poster", help="poster image (default: first frame)")
    p.add_argument("--normalize", action="store_true",
                   help="re-encode to H.264/AAC/yuv420p/CFR/faststart first")
    p.add_argument("--autoplay", action="store_true")
    args = p.parse_args()

    for tool in ("ffmpeg", "ffprobe"):
        if shutil.which(tool) is None:
            sys.exit(f"error: {tool} not found on PATH")

    workdir = Path(tempfile.mkdtemp(prefix="slide-video-"))
    video = args.video

    info = ffprobe(video)
    reasons = needs_normalization(info)
    if args.normalize and reasons:
        normalized = str(workdir / (Path(video).stem + "-h264.mp4"))
        print(f"normalizing ({'; '.join(reasons)}) ...")
        normalize(video, normalized)
        video = normalized
    elif reasons:
        print("warning: video may not play in PowerPoint: " + "; ".join(reasons))
        print("         re-run with --normalize to fix")

    # aspect-ratio sanity check
    v = next(s for s in ffprobe(video)["streams"] if s["codec_type"] == "video")
    var, box = int(v["width"]) / int(v["height"]), args.w / args.h
    if abs(var - box) / var > 0.02:
        print(f"warning: box aspect {box:.3f} != video aspect {var:.3f} "
              "(video will letterbox/stretch)")

    poster = args.poster
    if poster is None:
        poster = str(workdir / "poster.png")
        extract_poster(video, poster)

    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation(args.pptx)
    if not (1 <= args.slide <= len(prs.slides)):
        sys.exit(f"error: slide {args.slide} out of range (deck has {len(prs.slides)})")
    slide = prs.slides[args.slide - 1]

    # corruption guard (python-pptx issue #954)
    sld = slide._element
    for ac in sld.findall(f"{{{NS['mc']}}}AlternateContent"):
        if ac.find(f".//{{{NS['p']}}}timing") is not None:
            sys.exit(
                "error: this slide's <p:timing> is wrapped in mc:AlternateContent "
                "(non-standard animation present). add_movie() would duplicate the "
                "timing tree and corrupt the file. Embed the video on a fresh slide, "
                "or remove the exotic animation first."
            )

    movie = slide.shapes.add_movie(
        video, Inches(args.x), Inches(args.y), Inches(args.w), Inches(args.h),
        poster_frame_image=poster, mime_type="video/mp4",
    )
    cleaned_links = remove_empty_media_hyperlinks(movie)

    if args.autoplay:
        set_video_autoplay(sld, movie.shape_id)

    prs.save(args.pptx)
    mode = "autoplay" if args.autoplay else "click-to-play"
    print(f"ok: video embedded on slide {args.slide} ({mode}, shape id "
          f"{movie.shape_id}) in {args.pptx}")
    if cleaned_links:
        print(f"cleaned: removed {cleaned_links} empty media hyperlink reference(s)")


if __name__ == "__main__":
    main()
