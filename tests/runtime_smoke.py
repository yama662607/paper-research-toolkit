#!/usr/bin/env -S uv run --quiet
# /// script
# requires-python = ">=3.11"
# dependencies = ["python-pptx>=1.0", "pillow>=10.0"]
# ///
"""Cross-platform smoke test for slide-creator's read-only utility scripts."""

import json
import subprocess
import sys
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "skills" / "slide-creator" / "scripts"


def run(script: str, *args: object, expected: int = 0) -> subprocess.CompletedProcess:
    command = ["uv", "run", "--quiet", str(SCRIPTS / script), *map(str, args)]
    result = subprocess.run(command, capture_output=True, text=True, cwd=ROOT)
    if result.returncode != expected:
        raise AssertionError(
            f"{script} returned {result.returncode}, expected {expected}\n"
            f"stdout:\n{result.stdout}\nstderr:\n{result.stderr}"
        )
    return result


def build_fixture(path: Path, image_path: Path) -> None:
    Image.new("RGB", (32, 24), "#3A6B5C").save(image_path)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(5), Inches(0.6))
    title_run = box.text_frame.paragraphs[0].add_run()
    title_run.text = "Runtime fixture"
    title_run.font.size = Pt(18)
    source = slide.shapes.add_textbox(
        Inches(0.7), Inches(6.8), Inches(5), Inches(0.3)
    )
    source.name = "Source note"
    source_run = source.text_frame.paragraphs[0].add_run()
    source_run.text = "Source: runtime fixture"
    source_run.font.size = Pt(10)
    table = slide.shapes.add_table(
        1, 1, Inches(6), Inches(1.5), Inches(2), Inches(0.6)
    ).table
    table_run = table.cell(0, 0).text_frame.paragraphs[0].add_run()
    table_run.text = "Readable table"
    table_run.font.size = Pt(18)
    slide.shapes.add_picture(
        str(image_path), Inches(1), Inches(1.5), width=Inches(2), height=Inches(1.5)
    )
    prs.save(path)


def build_unknown_font_fixture(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    box.text_frame.paragraphs[0].text = "Theme-inherited size"
    prs.save(path)


def build_field_fixture(path: Path, size: str = "1000") -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.5))
    paragraph = box.text_frame.paragraphs[0]
    field = OxmlElement("a:fld")
    field.set("id", "{00000000-0000-0000-0000-000000000001}")
    field.set("type", "slidenum")
    rpr = OxmlElement("a:rPr")
    rpr.set("sz", size)
    field.append(rpr)
    text = OxmlElement("a:t")
    text.text = "1"
    field.append(text)
    paragraph._p.insert(0, field)
    prs.save(path)


def build_group_fixture(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    group.name = "Citation group"
    child = group.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.5))
    child.name = "Audience body"
    run = child.text_frame.paragraphs[0].add_run()
    run.text = "Small claim text"
    run.font.size = Pt(10)
    prs.save(path)


def main() -> None:
    with tempfile.TemporaryDirectory(prefix="slide-creator-runtime-") as tmp:
        root = Path(tmp)
        deck = root / "fixture.pptx"
        image = root / "fixture.png"
        build_fixture(deck, image)

        run("verify_deck.py", deck)
        run("verify_deck.py", deck, "--min-font-size", 16, expected=1)
        strict = run(
            "verify_deck.py",
            deck,
            "--min-font-size",
            16,
            "--allow-small-font-regex",
            r"^Source:",
        )
        assert "no non-exempt text below 16 pt" in strict.stdout
        strict_by_shape = run(
            "verify_deck.py",
            deck,
            "--min-font-size",
            16,
            "--allow-small-font-shape-regex",
            r"^Source note$",
        )
        assert "no non-exempt text below 16 pt" in strict_by_shape.stdout

        unknown_deck = root / "unknown-font.pptx"
        build_unknown_font_fixture(unknown_deck)
        unknown = run(
            "verify_deck.py",
            unknown_deck,
            "--min-font-size",
            16,
            expected=1,
        )
        assert "audience floor cannot be proven" in unknown.stdout

        field_deck = root / "field-font.pptx"
        build_field_fixture(field_deck)
        field_result = run(
            "verify_deck.py", field_deck, "--min-font-size", 16, expected=1
        )
        assert "10 pt" in field_result.stdout
        run(
            "verify_deck.py",
            field_deck,
            "--min-font-size",
            16,
            "--allow-small-font-regex",
            r"^1$",
        )

        nonfinite_deck = root / "nonfinite-font.pptx"
        build_field_fixture(nonfinite_deck, size="NaN")
        nonfinite = run(
            "verify_deck.py",
            nonfinite_deck,
            "--min-font-size",
            16,
            expected=1,
        )
        assert "audience floor cannot be proven" in nonfinite.stdout

        group_deck = root / "group-font.pptx"
        build_group_fixture(group_deck)
        parent_exemption = run(
            "verify_deck.py",
            group_deck,
            "--min-font-size",
            16,
            "--allow-small-font-shape-regex",
            r"^Citation group$",
            expected=1,
        )
        assert "Small claim text" in parent_exemption.stdout
        run(
            "verify_deck.py",
            group_deck,
            "--min-font-size",
            16,
            "--allow-small-font-shape-regex",
            r"^Audience body$",
        )

        invalid_floor = run(
            "verify_deck.py", deck, "--min-font-size", "nan", expected=2
        )
        assert "finite number greater than zero" in invalid_floor.stderr

        ingest = run("ingest_deck.py", deck, "--json")
        summary = json.loads(ingest.stdout)["summary"]
        assert summary["slide_count"] == 1

        layout = run("dump_layout.py", deck, "--slide", 1, "--json")
        assert json.loads(layout.stdout)["1"]

        out = root / "media"
        run("extract_media.py", deck, "--slide", 1, "--out", out)
        assert list(out.glob("media_1.*"))
        run("extract_media.py", deck, "--slide", 1, "--out", out, expected=2)
        run("extract_media.py", deck, "--slide", 1, "--out", out, "--overwrite")

        diff = run(
            "capture_edits.py",
            "--reference",
            deck,
            "--edited",
            deck,
            "--deep",
            "--json",
        )
        assert json.loads(diff.stdout)["changes"] == []

        report = root / "inspect.json"
        run("sync_from_pptx.py", deck, "--inspect-only", "--out", report)
        assert len(json.loads(report.read_text(encoding="utf-8"))["slides"]) == 1

        run("safe_rebuild.py", "--deck", deck, expected=2)
        backups = root / "backups"
        run(
            "safe_rebuild.py",
            "--deck",
            deck,
            "--reference",
            deck,
            "--backup-dir",
            backups,
        )
        assert len(list(backups.glob("fixture_backup_*.pptx"))) == 1

    print("slide-creator runtime smoke tests passed")


if __name__ == "__main__":
    main()
